from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "Ashburton Flood SAR Analysis Pipeline"
title_frame.paragraphs[0].font.size = Pt(28)
title_frame.paragraphs[0].font.bold = True

# Steps
steps = [
    ("Data Input\n- Pre/Post-flood SAR TIFF\n- Annotation XML", 0.5, 1.5, RGBColor(0, 102, 204)),
    ("Preprocessing: process_sar.py\n- Radiometric Calibration\n- Speckle Filtering\n- Ratio & Change Detection\n- Export GeoTIFF", 0.5, 3, RGBColor(0, 153, 0)),
    ("ML Data Preparation: prepare_ml_data.py\n- Normalize 2-channel SAR\n- Generate 256x256 patches", 0.5, 5, RGBColor(0, 153, 0)),
    ("U-Net Training: train_unet.py\n- Load patches & augmentation\n- Train 2-channel U-Net\n- Save model", 0.5, 7, RGBColor(255, 102, 0)),
    ("FastAPI Server: server/server.py\n- Ollama LLM Gateway API\n- PostGIS Integration\n- Serve Predictions", 0.5, 9, RGBColor(255, 102, 0)),
    ("Docker & DevContainer\n- Integrates B/C/D/E\n- Reproducible environment", 0.5, 11, RGBColor(128, 128, 128)),
]

for text, left, top, color in steps:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(9), Inches(1.5))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    for line in text.split("\n"):
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

prs.save("Ashburton_Flood_SAR_Pipeline.pptx")
print("PowerPoint saved!")

